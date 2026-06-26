"""
Gerador de Apresentacoes A&M
============================
Gera arquivos .pptx usando o Template_Base.pptx com os layouts e regras da A&M.

Uso:
    python3 gerar_apresentacao.py apresentacao.json
    python3 gerar_apresentacao.py  (usa exemplo embutido)

O JSON de entrada define os slides da apresentacao.
Veja a funcao _exemplo() para o formato completo.
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ============================================================================
# CONSTANTES - Paleta de cores A&M
# ============================================================================
CORES = {
    "marinho":      RGBColor(0x00, 0x2A, 0x46),
    "marinho_alt":  RGBColor(0x00, 0x2A, 0x47),
    "azul_escuro":  RGBColor(0x0E, 0x28, 0x41),
    "azul_claro":   RGBColor(0x5E, 0x8A, 0xB4),
    "branco":       RGBColor(0xFF, 0xFF, 0xFF),
    "preto":        RGBColor(0x00, 0x00, 0x00),
    "cinza_medio":  RGBColor(0xCC, 0xD5, 0xDB),
    "cinza_claro":  RGBColor(0xEB, 0xEE, 0xF1),
    "cinza_texto":  RGBColor(0x59, 0x59, 0x59),
    "cinza_inter":  RGBColor(0x33, 0x55, 0x6D),
    "laranja":      RGBColor(0xF7, 0x8C, 0x13),
    "vermelho":     RGBColor(0xC0, 0x00, 0x00),
    "verde":        RGBColor(0x00, 0xB0, 0x50),
}

# Cor de texto padrao por tema de fundo
COR_TEXTO_POR_TEMA = {
    "branco":      CORES["marinho"],
    "cinza_claro": CORES["marinho"],
    "cinza":       CORES["marinho"],
    "azul":        CORES["branco"],
    "marinho":     CORES["branco"],
}

# ============================================================================
# MAPA DE LAYOUTS - indice -> nome
# ============================================================================
LAYOUTS = {
    # 01 - Padrao Completo (tag + titulo + subtitulo + rodape)
    "padrao_branco":            0,
    "padrao_2linhas_branco":    1,
    "padrao_cinza_claro":       2,
    "padrao_cinza":             3,
    "padrao_azul":              4,
    "padrao_marinho":           5,

    # 02 - Padrao livre (so rodape)
    "livre_branco":             6,
    "livre_cinza_claro":        7,
    "livre_cinza":              8,
    "livre_azul":               9,
    "livre_marinho":            10,

    # 03 - Padrao Completo + Grafismo
    "grafismo_branco":          11,
    "grafismo_branco_2":        12,
    "grafismo_marinho":         13,
    "grafismo_marinho_2":       14,
    "grafismo_cinza":           15,
    "grafismo_cinza_claro":     16,

    # 04 - Destaque (frase de impacto)
    "destaque_branco":          17,
    "destaque_cinza":           18,
    "destaque_azul":            19,
    "destaque_marinho":         20,

    # 05 - Capa com Grafismo
    "capa_grafismo_branco":     21,
    "capa_grafismo_cinza":      22,
    "capa_grafismo_marinho":    23,

    # 06 - Capa com Foto
    "capa_foto_branco":         24,
    "capa_foto_cinza":          25,
    "capa_foto_marinho":        26,

    # 07 - Destaque com Foto
    "destaque_foto_branco":     27,
    "destaque_foto_cinza":      28,
    "destaque_foto_azul":       29,
    "destaque_foto_marinho":    30,

    # 08 - Agenda com Foto
    "agenda_cinza":             31,
    "agenda_azul":              32,
    "agenda_marinho":           33,

    # 09 - Divisoria de capitulo
    "divisoria_branco":         34,
    "divisoria_sub_branco":     35,
    "divisoria_cinza":          36,
    "divisoria_sub_cinza":      37,
    "divisoria_marinho":        38,
    "divisoria_sub_marinho":    39,

    # 10 - Completo + Grafismo (variacao)
    "completo_graf_branco":     40,
    "completo_graf_cinza_claro":41,
    "completo_graf_cinza":      42,
    "completo_graf_marinho":    43,

    # 11 - Texto 1 coluna
    "texto1_branco":            44,
    "texto1_diag_branco":       45,
    "texto1_marinho":           46,
    "texto1_diag_marinho":      47,

    # 12 - Texto 2 colunas
    "texto2_branco":            48,
    "texto2_diag_branco":       49,
    "texto2_marinho":           50,
    "texto2_diag_marinho":      51,

    # 13 - Texto 3 colunas
    "texto3_branco":            52,
    "texto3_diag_branco":       53,
    "texto3_marinho":           54,
    "texto3_diag_marinho":      55,

    # 13 - Texto 4 colunas
    "texto4_branco":            56,
    "texto4_diag_branco":       57,
    "texto4_marinho":           58,
    "texto4_diag_marinho":      59,

    # 14 - Texto 6 colunas
    "texto6_branco":            60,
    "texto6_diag_branco":       61,
    "texto6_marinho":           62,
    "texto6_diag_marinho":      63,

    # 15 - Texto 6 colunas com logo
    "texto6_logo_branco":       64,
    "texto6_logo_diag_branco":  65,
    "texto6_logo_marinho":      66,
    "texto6_logo_diag_marinho": 67,

    # 16 - Grafico
    "grafico_branco":           68,
    "grafico_cinza_claro":      69,
    "grafico_marinho":          70,

    # 17 - Diagonal
    "diagonal_branco":          71,
    "diagonal_azul":            72,
    "diagonal_marinho":         73,
    "diagonal_marinho_escuro":  74,

    # 18 - Topicos
    "topicos_claro":            75,
    "topicos_escuro":           76,

    # 19 - Equipe
    "equipe_6":                 77,
    "equipe_5":                 78,
    "equipe_4":                 79,
    "equipe_3":                 80,
    "equipe_2":                 81,
    "equipe_bio":               82,

    # 20 - Despedida
    "despedida":                83,

    # Blank
    "blank":                    84,
}

# ============================================================================
# MAPEAMENTO DE PLACEHOLDERS POR TIPO DE SLIDE
# ============================================================================
# Cada tipo mapeia nomes semanticos -> placeholder idx
# Variantes "branco" e "marinho" podem ter idx diferentes

PH_MAP = {
    # --- 01 Padrao Completo ---
    "padrao": {
        "branco":  {"titulo": 13, "tag": 14, "subtitulo": 15, "rodape": 16},
        "marinho": {"titulo": 13, "tag": 15, "subtitulo": 16, "rodape": 17},
        "cinza_claro": {"titulo": 13, "tag": 15, "subtitulo": 16, "rodape": 17},
        "cinza":   {"titulo": 13, "tag": 15, "subtitulo": 16, "rodape": 17},
        "azul":    {"titulo": 13, "tag": 15, "subtitulo": 16, "rodape": 17},
    },
    "padrao_2linhas": {
        "branco":  {"titulo": 13, "tag": 14, "rodape": 16},
    },

    # --- 03 Grafismo ---
    "grafismo": {
        "branco":  {"titulo": 13, "tag": 14, "subtitulo": 15, "rodape": 16},
        "marinho": {"titulo": 13, "tag": 15, "subtitulo": 16, "rodape": 17},
        "cinza":   {"titulo": 13, "tag": 15, "subtitulo": 16, "rodape": 17},
        "cinza_claro": {"titulo": 13, "tag": 15, "subtitulo": 16, "rodape": 17},
    },

    # --- 04 Destaque ---
    "destaque": {
        "_all": {"texto": 13, "credito": 14},
    },

    # --- 05 Capa Grafismo ---
    "capa_grafismo": {
        "_all": {"titulo": 13, "subtitulo": 14, "data": 15},
    },

    # --- 06 Capa Foto ---
    "capa_foto": {
        "_all": {"foto": 12, "titulo": 13, "subtitulo": 14, "data": 15},
    },

    # --- 07 Destaque Foto ---
    "destaque_foto": {
        "_all": {"foto": 12, "texto": 13},
    },

    # --- 08 Agenda ---
    "agenda": {
        "_all": {
            "foto": 12, "titulo": 13,
            "num1": 14, "item1": 15,
            "num2": 16, "item2": 17,
            "num3": 18, "item3": 19,
            "num4": 20, "item4": 21,
            "num5": 22, "item5": 23,
        },
    },

    # --- 09 Divisoria ---
    "divisoria": {
        "_all":     {"titulo": 12, "rodape": 16},
    },
    "divisoria_sub": {
        "_all":     {"titulo": 12, "subtitulo": 13, "rodape": 16},
    },

    # --- 11 Texto 1 coluna ---
    "texto1": {
        "branco":  {"titulo": 19, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "rodape": 12},
        "marinho": {"titulo": 16, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "rodape": 12},
    },

    # --- 12 Texto 2 colunas ---
    "texto2": {
        "branco":  {"titulo": 19, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 80, "col2_corpo": 81, "rodape": 12},
        "marinho": {"titulo": 16, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 80, "col2_corpo": 81, "rodape": 12},
    },

    # --- 13 Texto 3 colunas ---
    "texto3": {
        "branco":  {"titulo": 19, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 82, "col2_corpo": 83, "col3_titulo": 80, "col3_corpo": 81, "rodape": 12},
        "marinho": {"titulo": 16, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 82, "col2_corpo": 83, "col3_titulo": 80, "col3_corpo": 81, "rodape": 12},
    },

    # --- 13 Texto 4 colunas ---
    "texto4": {
        "branco":  {"titulo": 19, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 82, "col2_corpo": 83, "col3_titulo": 80, "col3_corpo": 81, "col4_titulo": 84, "col4_corpo": 85, "rodape": 12},
        "marinho": {"titulo": 19, "tag": 17, "subtitulo": 18, "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 82, "col2_corpo": 83, "col3_titulo": 80, "col3_corpo": 81, "col4_titulo": 84, "col4_corpo": 85, "rodape": 12},
    },

    # --- 14 Texto 6 colunas ---
    "texto6": {
        "branco":  {"titulo": 19, "tag": 17, "subtitulo": 18,
                    "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 82, "col2_corpo": 83,
                    "col3_titulo": 80, "col3_corpo": 81, "col4_titulo": 84, "col4_corpo": 85,
                    "col5_titulo": 88, "col5_corpo": 89, "col6_titulo": 86, "col6_corpo": 87,
                    "rodape": 12},
        "marinho": {"titulo": 16, "tag": 17, "subtitulo": 18,
                    "col1_titulo": 14, "col1_corpo": 15, "col2_titulo": 82, "col2_corpo": 83,
                    "col3_titulo": 80, "col3_corpo": 81, "col4_titulo": 84, "col4_corpo": 85,
                    "col5_titulo": 88, "col5_corpo": 89, "col6_titulo": 86, "col6_corpo": 87,
                    "rodape": 12},
    },

    # --- 16 Grafico ---
    "grafico": {
        "branco":  {"titulo": 23, "tag": 21, "subtitulo": 22,
                    "bloco1_titulo": 14, "bloco1_corpo": 15, "nota_grafico": 17,
                    "bloco2_titulo": 18, "bloco2_corpo": 19, "rodape": 12},
        "marinho": {"titulo": 16, "tag": 21, "subtitulo": 22,
                    "bloco1_titulo": 14, "bloco1_corpo": 15, "nota_grafico": 17,
                    "bloco2_titulo": 18, "bloco2_corpo": 19, "rodape": 12},
    },

    # --- 17 Diagonal ---
    "diagonal": {
        "branco":  {"titulo": 19, "tag": 17, "subtitulo": 18, "bloco_titulo": 13, "bloco_corpo": 14, "rodape": 15},
        "marinho": {"titulo": 16, "tag": 21, "subtitulo": 22, "bloco_titulo": 13, "bloco_corpo": 14, "rodape": 15},
    },

    # --- 18 Topicos ---
    "topicos": {
        "_all": {
            "titulo": 16, "tag": 21, "subtitulo": 22,
            "destaque": 35,
            "top1_titulo": 13, "top1_corpo": 15,
            "top2_titulo": 27, "top2_corpo": 28,
            "top3_titulo": 29, "top3_corpo": 30,
            "top4_titulo": 31, "top4_corpo": 32,
            "top5_titulo": 33, "top5_corpo": 34,
            "rodape": 36,
        },
    },
}


# ============================================================================
# FUNCOES DE FORMATACAO
# ============================================================================

def _set_text(placeholder, text, font_name=None, font_size=None, color=None,
              bold=None, italic=None, alignment=None):
    """Preenche um placeholder com texto formatado."""
    tf = placeholder.text_frame

    lines = text.split("\n") if text else [""]

    # Usar .text para a primeira linha (garante persistencia)
    tf.text = lines[0]

    # Formatar o paragrafo/run existente
    def _format_paragraph(p, line_text=None):
        if line_text is not None:
            p.text = line_text
        for run in p.runs:
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = Pt(font_size)
            if color:
                if isinstance(color, str) and color in CORES:
                    run.font.color.rgb = CORES[color]
                elif isinstance(color, RGBColor):
                    run.font.color.rgb = color
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
        if alignment:
            align_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }
            p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)

    _format_paragraph(tf.paragraphs[0])

    # Linhas adicionais
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        _format_paragraph(p)


def _set_text_multi(placeholder, runs_data):
    """Preenche um placeholder com multiplos runs formatados.

    runs_data: lista de dicts com keys: text, font_name, font_size, color, bold, italic
    """
    tf = placeholder.text_frame

    # Primeiro run: usar .text para garantir persistencia
    first = runs_data[0] if runs_data else {}
    tf.text = first.get("text", "")
    if tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        if first.get("font_name"):
            run.font.name = first["font_name"]
        if first.get("font_size"):
            run.font.size = Pt(first["font_size"])
        if first.get("color"):
            c = first["color"]
            if isinstance(c, str) and c in CORES:
                run.font.color.rgb = CORES[c]
        if first.get("bold"):
            run.font.bold = True
        if first.get("italic"):
            run.font.italic = True

    # Runs adicionais
    p = tf.paragraphs[0]
    for rd in runs_data[1:]:
        run = p.add_run()
        run.text = rd.get("text", "")
        if rd.get("font_name"):
            run.font.name = rd["font_name"]
        if rd.get("font_size"):
            run.font.size = Pt(rd["font_size"])
        if rd.get("color"):
            c = rd["color"]
            if isinstance(c, str) and c in CORES:
                run.font.color.rgb = CORES[c]
        if rd.get("bold"):
            run.font.bold = True
        if rd.get("italic"):
            run.font.italic = True


def _detectar_tema(layout_key):
    """Retorna o tema de cor baseado no nome do layout."""
    for tema in ["marinho", "azul", "cinza_claro", "cinza", "branco"]:
        if tema in layout_key:
            return tema
    return "branco"


def _get_ph_map(tipo_slide, tema):
    """Retorna o mapeamento de placeholders para o tipo e tema."""
    if tipo_slide not in PH_MAP:
        return {}
    mapa = PH_MAP[tipo_slide]
    if "_all" in mapa:
        return mapa["_all"]
    # Tentar tema exato, senao fallback
    if tema in mapa:
        return mapa[tema]
    # Fallback: primeiro disponivel
    return next(iter(mapa.values()))


# ============================================================================
# GERADOR PRINCIPAL
# ============================================================================

def _normalizar_cfg(cfg, mapa):
    """Normaliza nomes de campos do JSON (SKILL.md -> interno).

    Se o campo do SKILL.md existe no cfg mas o campo interno nao,
    copia o valor para o campo interno. Assim ambos os formatos funcionam.
    """
    out = dict(cfg)
    for skill_name, internal_name in mapa.items():
        if skill_name in out and internal_name not in out:
            out[internal_name] = out[skill_name]
    return out


# Mapas de traducao por tipo de elemento (SKILL.md -> interno)
_MAP_RETANGULO = {
    "x": "left", "y": "top", "largura": "width", "altura": "height",
    "cor_fundo": "preenchimento", "cor_borda": "borda",
    "largura_borda": "borda_espessura", "tamanho_fonte": "font_size",
    "negrito": "bold", "alinhamento": "alignment", "ancora_vertical": "vertical",
    "cor_texto": "cor_texto",  # mesmo nome, mantido por consistencia
}

_MAP_BADGE = {
    "x": "left", "y": "top", "largura": "width", "altura": "height",
    "tamanho_fonte": "font_size", "negrito": "bold",
}

_MAP_GRAFICO = {
    "x": "left", "y": "top", "largura": "width", "altura": "height",
    "mostrar_rotulos": "mostrar_valores",
}

_MAP_CALLOUT = {
    "x": "left", "y": "top", "largura": "width",
    "valor": "numero", "cor_valor": "cor_numero",
    "tamanho_valor": "font_size_numero", "tamanho_label": "font_size_label",
    "cor_fundo": "fundo",
}

_MAP_ICONE = {
    "x": "left", "y": "top", "tamanho": "size",
    "largura": "width", "altura": "height",
}

_MAP_SHAPE = {
    "x": "left", "y": "top", "largura": "width", "altura": "height",
    "texto": "text", "cor_texto": "color", "tamanho_fonte": "font_size",
    "negrito": "bold",
}

_MAP_LINHA = {
    "cor": "cor", "espessura": "espessura",
}


class GeradorAM:
    """Gerador de apresentacoes no padrao A&M."""

    def __init__(self, template_path="Template_Base.pptx"):
        self.template_path = template_path
        self.prs = None

    def gerar(self, config, output_path="Apresentacao.pptx"):
        """Gera a apresentacao a partir de um dict de configuracao.

        config = {
            "arquivo_saida": "Nome_Apresentacao.pptx",  (opcional)
            "slides": [
                { "tipo": "capa_grafismo_marinho", "titulo": "...", ... },
                ...
            ]
        }
        """
        self.prs = Presentation(self.template_path)
        output = config.get("arquivo_saida", output_path)

        warnings_total = []
        for i, slide_cfg in enumerate(config.get("slides", [])):
            try:
                cfg_efetivo = self._criar_slide(slide_cfg)
                w = self._validar_slide(i + 1, cfg_efetivo)
                warnings_total.extend(w)
            except Exception as e:
                print(f"  ERRO no slide {i+1}: {e}")

        self.prs.save(output)
        print(f"\nApresentacao salva: {output}")
        print(f"Total de slides: {len(self.prs.slides)}")

        if warnings_total:
            print(f"\n{'='*60}")
            print(f"AVISOS DE QUALIDADE ({len(warnings_total)} encontrados):")
            print(f"{'='*60}")
            for w in warnings_total:
                print(f"  {w}")
            print(f"{'='*60}")
            print("Revise os avisos acima e ajuste o JSON para melhorar a apresentacao.")

        return output

    # Mapa de equivalencias escuro -> claro para slides com graficos
    _LAYOUT_CLARO = {
        "padrao_marinho": "padrao_branco",
        "padrao_azul": "padrao_branco",
        "padrao_cinza": "padrao_branco",
        "grafismo_marinho": "grafismo_branco",
        "grafismo_cinza": "grafismo_branco",
        "completo_graf_marinho": "completo_graf_branco",
        "completo_graf_cinza": "completo_graf_branco",
        "texto1_marinho": "texto1_branco",
        "texto2_marinho": "texto2_branco",
        "texto3_marinho": "texto3_branco",
        "texto4_marinho": "texto4_branco",
        "texto6_marinho": "texto6_branco",
        "grafico_marinho": "grafico_branco",
        "livre_marinho": "livre_branco",
        "livre_azul": "livre_branco",
    }

    def _criar_slide(self, cfg):
        """Cria um slide a partir da configuracao."""
        tipo = cfg.get("tipo", "")

        # Regra: se o slide tem graficos, forcar variante clara
        if cfg.get("graficos") and tipo in self._LAYOUT_CLARO:
            tipo_original = tipo
            tipo = self._LAYOUT_CLARO[tipo]
            cfg = dict(cfg, tipo=tipo)
            print(f"  [auto] Layout '{tipo_original}' -> '{tipo}' (slide com grafico)")

        layout_idx = LAYOUTS.get(tipo)

        if layout_idx is None:
            # Tentar como indice numerico direto
            if isinstance(tipo, int):
                layout_idx = tipo
            else:
                print(f"  Layout '{tipo}' nao encontrado. Usando blank.")
                layout_idx = 84

        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)

        # Identificar tipo base e tema para buscar mapeamento
        tipo_base = self._extrair_tipo_base(tipo)
        tema = _detectar_tema(tipo)
        ph_map = _get_ph_map(tipo_base, tema)

        # Montar dict de placeholders disponiveis
        ph_dict = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        # Preencher cada campo que esta no cfg e tem placeholder mapeado
        for campo, idx in ph_map.items():
            valor = cfg.get(campo)
            if valor is None:
                continue
            if idx not in ph_dict:
                continue

            ph = ph_dict[idx]

            # Se o campo e uma foto, pular (tratado separadamente)
            if campo == "foto":
                if os.path.exists(str(valor)):
                    ph.insert_picture(valor)
                continue

            # Texto simples
            if isinstance(valor, str):
                self._aplicar_texto_padrao(ph, campo, valor, tema)
            # Lista de runs
            elif isinstance(valor, list):
                _set_text_multi(ph, valor)

        # Campos especiais para agenda (itens numerados)
        if tipo_base == "agenda":
            self._preencher_agenda(slide, cfg, ph_map)

        # Campos especiais para topicos
        if tipo_base == "topicos":
            self._preencher_topicos(slide, cfg, ph_map)

        # Campos de colunas (texto2, texto3, texto4, texto6)
        if tipo_base.startswith("texto") and tipo_base != "texto1":
            self._preencher_colunas(slide, cfg, ph_map, tipo)

        # Tabela (qualquer slide pode ter uma)
        if "tabela" in cfg:
            self._adicionar_tabela(slide, cfg["tabela"], tema)

        # Graficos
        for graf_cfg in cfg.get("graficos", []):
            self._adicionar_grafico(slide, graf_cfg, tema)

        # Retangulos / Cards
        for ret_cfg in cfg.get("retangulos", []):
            self._adicionar_retangulo(slide, ret_cfg, tema)

        # Badges (pills)
        for badge_cfg in cfg.get("badges", []):
            self._adicionar_badge(slide, badge_cfg, tema)

        # Linhas decorativas
        for linha_cfg in cfg.get("linhas", []):
            self._adicionar_linha(slide, linha_cfg, tema)

        # Icones circulares
        for icone_cfg in cfg.get("icones", []):
            self._adicionar_icone(slide, icone_cfg, tema)

        # Callouts de KPI (numero grande + label)
        for callout_cfg in cfg.get("callouts", []):
            self._adicionar_callout(slide, callout_cfg, tema)

        # Shapes de texto livre (posicionamento manual)
        for shape_cfg in cfg.get("shapes", []):
            self._adicionar_shape(slide, shape_cfg, tema)

        # Componentes DD (extensao A&M)
        for wf_cfg in cfg.get("waterfalls", []):
            self._adicionar_waterfall(slide, wf_cfg, tema)
        for q_cfg in cfg.get("quadrantes", []):
            self._adicionar_quadrante(slide, q_cfg, tema)
        for f_cfg in cfg.get("funis", []):
            self._adicionar_funil(slide, f_cfg, tema)

        return cfg

    # Tipos de slide que nao precisam de elementos visuais extras
    _TIPOS_SEM_VISUAL = {"capa_grafismo", "capa_foto", "destaque", "destaque_foto",
                         "divisoria", "divisoria_sub", "agenda", "despedida", "blank",
                         "equipe", "topicos"}

    def _validar_slide(self, num, cfg):
        """Valida qualidade do slide e retorna lista de warnings."""
        warnings = []
        tipo = cfg.get("tipo", "")
        tipo_base = self._extrair_tipo_base(tipo)

        # Slides que nao precisam de validacao de elementos visuais
        if tipo_base in self._TIPOS_SEM_VISUAL:
            return warnings

        # 1. Verificar se tem ao menos 1 elemento visual
        tem_retangulos = bool(cfg.get("retangulos"))
        tem_graficos = bool(cfg.get("graficos"))
        tem_callouts = bool(cfg.get("callouts"))
        tem_badges = bool(cfg.get("badges"))
        tem_icones = bool(cfg.get("icones"))
        tem_tabela = bool(cfg.get("tabela"))
        tem_colunas = bool(cfg.get("colunas"))
        tem_dd = bool(cfg.get("waterfalls") or cfg.get("quadrantes") or cfg.get("funis"))

        elementos_visuais = sum([tem_retangulos, tem_graficos, tem_callouts,
                                 tem_badges, tem_tabela, tem_colunas, tem_dd])

        if elementos_visuais == 0:
            warnings.append(f"Slide {num} ({tipo}): SEM ELEMENTOS VISUAIS - adicione retangulos, graficos, callouts ou tabela")

        # 2. Verificar se tem icones (todo slide de conteudo deveria ter)
        if not tem_icones and not tem_colunas:
            warnings.append(f"Slide {num} ({tipo}): SEM ICONES - adicione icones para enriquecer visualmente")

        # 3. Verificar conteudo dos retangulos (muito pouco texto = slide pobre)
        for i, ret in enumerate(cfg.get("retangulos", [])):
            texto = ret.get("texto", "")
            if texto and len(texto) < 60:
                warnings.append(f"Slide {num} ({tipo}): retangulo {i+1} com pouco conteudo ({len(texto)} chars) - considere adicionar mais detalhes")
            font_size = ret.get("tamanho_fonte", ret.get("font_size", 12))
            if font_size < 10:
                warnings.append(f"Slide {num} ({tipo}): retangulo {i+1} com fonte {font_size}pt < 10pt minimo")

        # 4. Verificar graficos em fundo escuro (deve ter sido auto-corrigido)
        tema = _detectar_tema(tipo)
        if tem_graficos and tema in ("marinho", "azul"):
            warnings.append(f"Slide {num} ({tipo}): GRAFICO EM FUNDO ESCURO - use variante _branco")

        # 5. Verificar callouts sem retangulos/contexto
        if tem_callouts and not tem_retangulos and not tem_graficos and not tem_colunas:
            warnings.append(f"Slide {num} ({tipo}): callouts sozinhos sem conteudo - adicione retangulos ou graficos")

        return warnings

    def _extrair_tipo_base(self, tipo):
        """Extrai o tipo base do nome do layout."""
        # Mapeamento de layout key -> tipo base no PH_MAP
        mapeamentos = {
            "padrao": "padrao",
            "padrao_2linhas": "padrao_2linhas",
            "grafismo": "grafismo",
            "completo_graf": "grafismo",
            "destaque_foto": "destaque_foto",
            "destaque": "destaque",
            "capa_grafismo": "capa_grafismo",
            "capa_foto": "capa_foto",
            "agenda": "agenda",
            "divisoria_sub": "divisoria_sub",
            "divisoria": "divisoria",
            "texto1": "texto1",
            "texto2": "texto2",
            "texto3": "texto3",
            "texto4": "texto4",
            "texto6": "texto6",
            "grafico": "grafico",
            "diagonal": "diagonal",
            "topicos": "topicos",
            "livre": "livre",
            "equipe": "equipe",
            "despedida": "despedida",
            "blank": "blank",
        }
        # Tenta match do mais especifico pro mais generico
        for key in sorted(mapeamentos.keys(), key=len, reverse=True):
            if tipo.startswith(key) or key in tipo:
                return mapeamentos[key]
        return tipo

    def _aplicar_texto_padrao(self, ph, campo, texto, tema):
        """Aplica texto com formatacao padrao baseada no campo."""
        cor = COR_TEXTO_POR_TEMA.get(tema, CORES["marinho"])
        cor_sub = CORES["cinza_texto"] if tema in ("branco", "cinza_claro") else CORES["cinza_medio"]

        # Definir formatacao por campo
        formatos = {
            "titulo":       {"font": "Arial Nova Light",          "size": 27, "color": cor, "bold": False},
            "tag":          {"font": "Arial Nova Condensed",      "size": 12, "color": cor_sub, "bold": False},
            "subtitulo":    {"font": "Arial Nova Condensed",      "size": 16, "color": cor, "bold": True},
            "rodape":       {"font": "Arial Nova Condensed",      "size": 9,  "color": cor_sub, "bold": False},
            "texto":        {"font": "Arial Nova Light",          "size": 27, "color": cor, "bold": False},
            "credito":      {"font": "Arial Nova Condensed",      "size": 12, "color": cor_sub, "bold": False},
            "data":         {"font": "Arial Nova Condensed",      "size": 12, "color": cor_sub, "bold": False},
            "destaque":     {"font": "Arial Nova Light",          "size": 22, "color": cor, "bold": False},
            "nota_grafico": {"font": "Arial Nova Condensed",      "size": 9,  "color": cor_sub, "bold": False},
        }

        # Campos de coluna
        if "titulo" in campo and campo != "titulo":
            fmt = {"font": "Arial Nova Light", "size": 16, "color": cor, "bold": False}
        elif "corpo" in campo:
            fmt = {"font": "Arial Nova Condensed", "size": 12, "color": cor, "bold": False}
        else:
            fmt = formatos.get(campo, {"font": "Arial Nova Condensed", "size": 14, "color": cor, "bold": False})

        _set_text(ph, texto,
                  font_name=fmt["font"],
                  font_size=fmt["size"],
                  color=fmt["color"],
                  bold=fmt["bold"])

    def _preencher_agenda(self, slide, cfg, ph_map):
        """Preenche itens numerados da agenda."""
        ph_dict = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        itens = cfg.get("itens", [])
        for i, item in enumerate(itens):
            if i >= 5:
                break
            num_key = f"num{i+1}"
            item_key = f"item{i+1}"
            num_idx = ph_map.get(num_key)
            item_idx = ph_map.get(item_key)
            if num_idx is not None and num_idx in ph_dict:
                _set_text(ph_dict[num_idx], str(i + 1),
                          font_name="Arial Nova Light", font_size=27,
                          color=CORES["laranja"], bold=False)
            if item_idx is not None and item_idx in ph_dict:
                _set_text(ph_dict[item_idx], item,
                          font_name="Arial Nova Condensed", font_size=16,
                          bold=True)

    def _preencher_topicos(self, slide, cfg, ph_map):
        """Preenche topicos numerados."""
        ph_dict = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        topicos = cfg.get("topicos", [])
        n_topicos = min(len(topicos), 5)

        # Ajustar tamanho de fonte baseado no numero de topicos para evitar sobreposicao
        if n_topicos <= 3:
            titulo_size = 18
            corpo_size = 12
        elif n_topicos == 4:
            titulo_size = 15
            corpo_size = 10
        else:
            titulo_size = 13
            corpo_size = 9

        for i, topico in enumerate(topicos):
            if i >= 5:
                break
            titulo_key = f"top{i+1}_titulo"
            corpo_key = f"top{i+1}_corpo"
            if isinstance(topico, dict):
                t_titulo = topico.get("titulo", "")
                t_corpo = topico.get("corpo", "")
            else:
                t_titulo = str(i + 1)
                t_corpo = str(topico)

            titulo_idx = ph_map.get(titulo_key)
            corpo_idx = ph_map.get(corpo_key)
            if titulo_idx is not None and titulo_idx in ph_dict:
                ph = ph_dict[titulo_idx]
                _set_text(ph, t_titulo,
                          font_name="Arial Nova Light", font_size=titulo_size,
                          color=CORES["laranja"])
                # Habilitar auto-shrink para evitar overflow
                bodyPr = ph.text_frame._txBody.find(qn('a:bodyPr'))
                if bodyPr is not None:
                    bodyPr.set('wrap', 'square')
                    # Remover normAutofit ou spAutoFit existente
                    for child in list(bodyPr):
                        if child.tag.endswith('AutoFit') or child.tag.endswith('autoFit'):
                            bodyPr.remove(child)

            if corpo_idx is not None and corpo_idx in ph_dict:
                ph = ph_dict[corpo_idx]
                _set_text(ph, t_corpo,
                          font_name="Arial Nova Condensed", font_size=corpo_size)
                # Habilitar auto-shrink para evitar overflow
                bodyPr = ph.text_frame._txBody.find(qn('a:bodyPr'))
                if bodyPr is not None:
                    bodyPr.set('wrap', 'square')
                    for child in list(bodyPr):
                        if child.tag.endswith('AutoFit') or child.tag.endswith('autoFit'):
                            bodyPr.remove(child)

    def _preencher_colunas(self, slide, cfg, ph_map, tipo_layout):
        """Preenche colunas para layouts multi-coluna."""
        ph_dict = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        tema = _detectar_tema(tipo_layout)
        colunas = cfg.get("colunas", [])
        for i, col in enumerate(colunas):
            titulo_key = f"col{i+1}_titulo"
            corpo_key = f"col{i+1}_corpo"
            if isinstance(col, dict):
                c_titulo = col.get("titulo", "")
                c_corpo = col.get("corpo", "")
            else:
                c_titulo = ""
                c_corpo = str(col)

            titulo_idx = ph_map.get(titulo_key)
            corpo_idx = ph_map.get(corpo_key)
            if titulo_idx is not None and titulo_idx in ph_dict:
                self._aplicar_texto_padrao(
                    ph_dict[titulo_idx],
                    titulo_key, c_titulo, tema)
            if corpo_idx is not None and corpo_idx in ph_dict:
                self._aplicar_texto_padrao(
                    ph_dict[corpo_idx],
                    corpo_key, c_corpo, tema)


    def _adicionar_tabela(self, slide, tabela_cfg, tema):
        """Adiciona uma tabela ao slide.

        tabela_cfg = {
            "left": 0.5,   # polegadas
            "top": 1.8,
            "width": 12.3,
            "height": 4.5,
            "cabecalho": ["Col1", "Col2", "Col3"],
            "linhas": [
                ["val1", "val2", "val3"],
                ["val1", "val2", "val3"],
            ],
            "estilo": "claro" | "escuro"  (opcional)
        }
        """
        left = Inches(tabela_cfg.get("left", 0.5))
        top = Inches(tabela_cfg.get("top", 1.8))
        width = Inches(tabela_cfg.get("width", 12.3))
        height = Inches(tabela_cfg.get("height", 4.0))

        cabecalho = tabela_cfg.get("cabecalho", [])
        linhas = tabela_cfg.get("linhas", [])
        estilo = tabela_cfg.get("estilo", "claro")

        n_rows = len(linhas) + (1 if cabecalho else 0)
        n_cols = len(cabecalho) if cabecalho else (len(linhas[0]) if linhas else 1)

        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        # Cores baseadas no tema
        if estilo == "escuro":
            header_bg = CORES["marinho"]
            header_fg = CORES["branco"]
            row_bg_1 = CORES["cinza_claro"]
            row_bg_2 = RGBColor(0xFF, 0xFF, 0xFF)
            cell_fg = CORES["marinho"]
        else:
            header_bg = CORES["marinho"]
            header_fg = CORES["branco"]
            row_bg_1 = RGBColor(0xFF, 0xFF, 0xFF)
            row_bg_2 = CORES["cinza_claro"]
            cell_fg = CORES["marinho"]

        def _format_cell(cell, text, font_name, font_size, color, bold=False, bg=None):
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(text)
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
            cell.text_frame.word_wrap = True
            # Margem interna
            cell.text_frame.margin_left = Inches(0.08)
            cell.text_frame.margin_right = Inches(0.08)
            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_bottom = Inches(0.05)
            if bg:
                from pptx.oxml.ns import qn
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': f'{bg}'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

        def _rgb_to_hex(rgb_color):
            return f'{rgb_color}'

        row_offset = 0
        if cabecalho:
            for j, header_text in enumerate(cabecalho):
                _format_cell(table.cell(0, j), header_text,
                             "Arial Nova Condensed", 10, header_fg, bold=True,
                             bg=_rgb_to_hex(header_bg))
            row_offset = 1

        for i, linha in enumerate(linhas):
            bg = row_bg_1 if i % 2 == 0 else row_bg_2
            for j, valor in enumerate(linha):
                if j < n_cols:
                    _format_cell(table.cell(i + row_offset, j), valor,
                                 "Arial Nova Condensed", 9, cell_fg, bold=False,
                                 bg=_rgb_to_hex(bg))

    # ==================================================================
    # ELEMENTOS VISUAIS AVANCADOS
    # ==================================================================

    def _resolver_cor(self, cor_ref):
        """Resolve uma referencia de cor (string ou hex) para RGBColor."""
        if isinstance(cor_ref, RGBColor):
            return cor_ref
        if isinstance(cor_ref, str):
            if cor_ref in CORES:
                return CORES[cor_ref]
            # Hex direto: "#002B49" ou "002B49"
            hex_str = cor_ref.lstrip("#")
            if len(hex_str) == 6:
                return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        return CORES["marinho"]

    def _set_shape_fill(self, shape, cor_ref):
        """Define preenchimento solido de um shape."""
        cor = self._resolver_cor(cor_ref)
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor

    def _set_shape_border(self, shape, cor_ref, espessura=1.0):
        """Define borda de um shape."""
        cor = self._resolver_cor(cor_ref)
        shape.line.color.rgb = cor
        shape.line.width = Pt(espessura)

    def _set_shape_text(self, shape, texto, font_name="Arial Nova Condensed",
                        font_size=12, color=None, bold=False, italic=False,
                        alignment="center", vertical="middle"):
        """Configura texto dentro de um shape."""
        tf = shape.text_frame
        tf.word_wrap = True
        # Margem interna reduzida para melhor aproveitamento
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

        # Ancoragem vertical
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

        try:
            tf.auto_size = None
        except:
            pass

        # Vertical anchor via XML
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            anchor_val = {"top": "t", "middle": "ctr", "bottom": "b"}.get(vertical, "ctr")
            bodyPr.set('anchor', anchor_val)

        cor = self._resolver_cor(color) if color else CORES["marinho"]

        lines = texto.split("\n") if texto else [""]
        tf.text = lines[0]

        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

        def _fmt_para(p):
            p.alignment = align_map.get(alignment, PP_ALIGN.CENTER)
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = cor
                run.font.bold = bold
                run.font.italic = italic

        _fmt_para(tf.paragraphs[0])
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            _fmt_para(p)

    def _adicionar_grafico(self, slide, graf_cfg_raw, tema):
        """Adiciona um grafico ao slide.

        Aceita campos no formato SKILL.md (x, y, largura, etc.)
        ou no formato interno (left, top, width, etc.).
        """
        graf_cfg = _normalizar_cfg(graf_cfg_raw, _MAP_GRAFICO)

        left = Inches(graf_cfg.get("left", 0.5))
        top = Inches(graf_cfg.get("top", 2.0))
        width = Inches(graf_cfg.get("width", 5.0))
        height = Inches(graf_cfg.get("height", 3.5))

        tipo = graf_cfg.get("tipo", "barra")
        tipo_map = {
            "barra": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "barra_agrupada": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "barra_empilhada": XL_CHART_TYPE.COLUMN_STACKED,
            "barra_empilhada_100": XL_CHART_TYPE.COLUMN_STACKED_100,
            "barra_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
            "linha": XL_CHART_TYPE.LINE_MARKERS,
            "pizza": XL_CHART_TYPE.PIE,
            "rosca": XL_CHART_TYPE.DOUGHNUT,
        }
        chart_type = tipo_map.get(tipo, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()
        chart_data.categories = graf_cfg.get("categorias", [])
        for serie in graf_cfg.get("series", []):
            chart_data.add_series(serie.get("nome", ""), serie.get("valores", []))

        chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
        chart = chart_frame.chart

        # Titulo
        if graf_cfg.get("titulo"):
            chart.has_title = True
            chart.chart_title.text_frame.text = graf_cfg["titulo"]
            for p in chart.chart_title.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Arial Nova Condensed"
                    run.font.size = Pt(11)
                    run.font.color.rgb = CORES["marinho"]
                    run.font.bold = True
        else:
            chart.has_title = False

        # Legenda
        if graf_cfg.get("mostrar_legenda", True) and len(graf_cfg.get("series", [])) > 1:
            chart.has_legend = True
            pos_map = {"bottom": XL_LEGEND_POSITION.BOTTOM, "right": XL_LEGEND_POSITION.RIGHT,
                       "left": XL_LEGEND_POSITION.LEFT, "top": XL_LEGEND_POSITION.TOP}
            chart.legend.position = pos_map.get(graf_cfg.get("posicao_legenda", "bottom"), XL_LEGEND_POSITION.BOTTOM)
            chart.legend.include_in_layout = False
            chart.legend.font.name = "Arial Nova Condensed"
            chart.legend.font.size = Pt(9)
        else:
            chart.has_legend = False

        # Cores das series / fatias
        series_list = graf_cfg.get("series", [])
        if tipo in ("pizza", "rosca"):
            # Para pizza/rosca: colorir cada fatia (data point) individualmente
            categorias = graf_cfg.get("categorias", [])
            # Paleta de cores para fatias
            cores_fatia = graf_cfg.get("cores_fatias", None)
            if cores_fatia is None:
                # Usar paleta padrao A&M com cores distintas
                cores_fatia = ["#002B49", "#5E8AB4", "#F78C13", "#CCD5DB", "#33556D",
                               "#EBEEF1", "#0E2841", "#C00000", "#00B050", "#595959"]
            if len(chart.series) > 0:
                series = chart.series[0]
                for j in range(len(categorias)):
                    pt = series.points[j]
                    cor = self._resolver_cor(cores_fatia[j % len(cores_fatia)])
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = cor
        else:
            for i, serie_cfg in enumerate(series_list):
                if i < len(chart.series):
                    cor = self._resolver_cor(serie_cfg.get("cor", "marinho"))
                    chart.series[i].format.fill.solid()
                    chart.series[i].format.fill.fore_color.rgb = cor

        # Data labels
        mostrar = graf_cfg.get("mostrar_valores", False)
        # Para pizza/rosca, sempre mostrar labels com nome da categoria
        if tipo in ("pizza", "rosca"):
            mostrar = True
        if mostrar:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.name = "Arial Nova Condensed"
            data_labels.font.size = Pt(9)
            if tipo in ("pizza", "rosca"):
                # Mostrar nome da categoria + percentual (global)
                data_labels.show_category_name = True
                data_labels.show_percentage = True
                data_labels.show_value = False
                data_labels.show_series_name = False
                data_labels.font.color.rgb = CORES["marinho"]
                data_labels.number_format = '0%'
                try:
                    data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
                except:
                    pass
                # Cor de cada label baseada na luminosidade da fatia
                if len(chart.series) > 0:
                    from lxml import etree
                    series = chart.series[0]
                    cores_usadas = cores_fatia if cores_fatia is not None else []
                    for j in range(len(categorias)):
                        try:
                            pt = series.points[j]
                            cor_hex = cores_usadas[j % len(cores_usadas)] if cores_usadas else "#002B49"
                            cor_hex = cor_hex.lstrip("#")
                            r, g, b = int(cor_hex[0:2], 16), int(cor_hex[2:4], 16), int(cor_hex[4:6], 16)
                            luminosidade = (0.299 * r + 0.587 * g + 0.114 * b)
                            cor_label = CORES["branco"] if luminosidade < 140 else CORES["marinho"]
                            # Primeiro setar a cor (isso cria o dLbl customizado)
                            pt.data_label.font.color.rgb = cor_label
                            pt.data_label.font.name = "Arial Nova Condensed"
                            pt.data_label.font.size = Pt(9)
                            # Agora o dLbl existe - garantir que mostra categoria e percentual
                            dLbl = pt.data_label._dLbl
                            if dLbl is not None:
                                for tag_name, val in [('c:showCatName', '1'), ('c:showPercent', '1'),
                                                      ('c:showVal', '0'), ('c:showSerName', '0')]:
                                    elem = dLbl.find(qn(tag_name))
                                    if elem is None:
                                        elem = etree.SubElement(dLbl, qn(tag_name))
                                    elem.set('val', val)
                        except Exception:
                            pass
            else:
                data_labels.font.color.rgb = CORES["marinho"]
                data_labels.number_format = graf_cfg.get("formato_valores", "0")
                if tipo in ("barra_empilhada", "barra_empilhada_100"):
                    data_labels.label_position = XL_LABEL_POSITION.CENTER
                else:
                    data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END

        # Estilo dos eixos
        if chart_type not in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
            try:
                cat_axis = chart.category_axis
                cat_axis.tick_labels.font.name = "Arial Nova Condensed"
                cat_axis.tick_labels.font.size = Pt(9)
                cat_axis.tick_labels.font.color.rgb = CORES["cinza_texto"]
                cat_axis.has_major_gridlines = False

                val_axis = chart.value_axis
                val_axis.tick_labels.font.name = "Arial Nova Condensed"
                val_axis.tick_labels.font.size = Pt(9)
                val_axis.tick_labels.font.color.rgb = CORES["cinza_texto"]
                val_axis.has_major_gridlines = True
                val_axis.major_gridlines.format.line.color.rgb = CORES["cinza_claro"]
                val_axis.major_gridlines.format.line.width = Pt(0.5)

                if graf_cfg.get("eixo_y_titulo"):
                    val_axis.has_title = True
                    val_axis.axis_title.text_frame.text = graf_cfg["eixo_y_titulo"]
                    for p in val_axis.axis_title.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.name = "Arial Nova Condensed"
                            run.font.size = Pt(9)
            except Exception:
                pass

        # Fundo transparente
        try:
            chart.chart_style = 2
            plot_area = chart.plots[0]
            chart_frame.chart.element.find('.//' + qn('c:plotArea')).find(qn('c:spPr'))
        except Exception:
            pass

    def _adicionar_retangulo(self, slide, ret_cfg_raw, tema):
        """Adiciona um retangulo/card ao slide.

        Aceita campos no formato SKILL.md (x, y, largura, altura, cor_fundo, etc.)
        ou no formato interno (left, top, width, height, preenchimento, etc.).
        """
        ret_cfg = _normalizar_cfg(ret_cfg_raw, _MAP_RETANGULO)

        # raio_canto (SKILL.md) -> arredondado (interno): qualquer valor > 0 = True
        if "raio_canto" in ret_cfg and "arredondado" not in ret_cfg:
            ret_cfg["arredondado"] = bool(ret_cfg["raio_canto"])

        left = Inches(ret_cfg.get("left", 0.5))
        top = Inches(ret_cfg.get("top", 2.0))
        width = Inches(ret_cfg.get("width", 5.0))
        height = Inches(ret_cfg.get("height", 2.0))

        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if ret_cfg.get("arredondado", False) else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        # Preenchimento
        preencher = ret_cfg.get("preenchimento", "nenhum")
        if preencher and preencher != "nenhum":
            self._set_shape_fill(shape, preencher)
        else:
            shape.fill.background()

        # Borda
        borda = ret_cfg.get("borda", "nenhum")
        if borda and borda != "nenhum":
            self._set_shape_border(shape, borda, ret_cfg.get("borda_espessura", 1.0))
        else:
            shape.line.fill.background()

        # Texto
        if ret_cfg.get("texto"):
            self._set_shape_text(
                shape, ret_cfg["texto"],
                font_name=ret_cfg.get("font_name", "Arial Nova Condensed"),
                font_size=ret_cfg.get("font_size", 12),
                color=ret_cfg.get("cor_texto", "marinho"),
                bold=ret_cfg.get("bold", False),
                italic=ret_cfg.get("italic", False),
                alignment=ret_cfg.get("alignment", "left"),
                vertical=ret_cfg.get("vertical", "top"),
            )

    def _adicionar_badge(self, slide, badge_cfg_raw, tema):
        """Adiciona um badge (pill) ao slide.

        Aceita campos no formato SKILL.md (x, y, largura, etc.)
        ou no formato interno (left, top, width, etc.).
        """
        badge_cfg = _normalizar_cfg(badge_cfg_raw, _MAP_BADGE)

        left = Inches(badge_cfg.get("left", 0.5))
        top = Inches(badge_cfg.get("top", 0.5))
        width = Inches(badge_cfg.get("width", 2.0))
        height = Inches(badge_cfg.get("height", 0.35))

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

        # Fundo
        cor_fundo = badge_cfg.get("cor_fundo", "marinho")
        if cor_fundo and cor_fundo != "nenhum":
            self._set_shape_fill(shape, cor_fundo)
        else:
            shape.fill.background()

        # Borda
        borda = badge_cfg.get("borda", "nenhum")
        if borda and borda != "nenhum":
            self._set_shape_border(shape, badge_cfg.get("borda_cor", borda), 1.0)
        else:
            shape.line.fill.background()

        # Texto
        self._set_shape_text(
            shape, badge_cfg.get("texto", ""),
            font_name=badge_cfg.get("font_name", "Arial Nova Condensed"),
            font_size=badge_cfg.get("font_size", 10),
            color=badge_cfg.get("cor_texto", "branco"),
            bold=badge_cfg.get("bold", True),
            alignment="center",
            vertical="middle",
        )

    def _adicionar_linha(self, slide, linha_cfg_raw, tema):
        """Adiciona uma linha decorativa ao slide.

        Aceita formato SKILL.md (x, y, comprimento, direcao)
        ou formato interno (x1, y1, x2, y2).
        """
        linha_cfg = _normalizar_cfg(linha_cfg_raw, _MAP_LINHA)

        # Converter formato SKILL.md (x, y, comprimento, direcao) -> (x1, y1, x2, y2)
        if "x" in linha_cfg and "x1" not in linha_cfg:
            x = linha_cfg["x"]
            y = linha_cfg["y"]
            comp = linha_cfg.get("comprimento", 12.0)
            direcao = linha_cfg.get("direcao", "horizontal")
            if direcao == "vertical":
                linha_cfg["x1"] = x
                linha_cfg["y1"] = y
                linha_cfg["x2"] = x
                linha_cfg["y2"] = y + comp
            else:  # horizontal (default)
                linha_cfg["x1"] = x
                linha_cfg["y1"] = y
                linha_cfg["x2"] = x + comp
                linha_cfg["y2"] = y

        x1 = Inches(linha_cfg.get("x1", 0.5))
        y1 = Inches(linha_cfg.get("y1", 4.5))
        x2 = Inches(linha_cfg.get("x2", 12.5))
        y2 = Inches(linha_cfg.get("y2", 4.5))

        connector = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            x1, y1, x2, y2
        )

        cor = self._resolver_cor(linha_cfg.get("cor", "laranja"))
        connector.line.color.rgb = cor
        connector.line.width = Pt(linha_cfg.get("espessura", 2.0))

        if linha_cfg.get("tracejado", False):
            # Dash style via XML
            ln = connector.line._ln
            ln.set('dash', 'dash')

    # Mapa de formas nativas do PowerPoint para icones
    SHAPE_MAP = {
        "circulo": MSO_SHAPE.OVAL,
        "engrenagem": MSO_SHAPE.GEAR_6,
        "engrenagem9": MSO_SHAPE.GEAR_9,
        "raio": MSO_SHAPE.LIGHTNING_BOLT,
        "estrela": MSO_SHAPE.STAR_5_POINT,
        "estrela6": MSO_SHAPE.STAR_6_POINT,
        "diamante": MSO_SHAPE.DIAMOND,
        "pentagono": MSO_SHAPE.PENTAGON,
        "hexagono": MSO_SHAPE.HEXAGON,
        "octagono": MSO_SHAPE.OCTAGON,
        "cubo": MSO_SHAPE.CUBE,
        "nuvem": MSO_SHAPE.CLOUD,
        "sol": MSO_SHAPE.SUN,
        "coracao": MSO_SHAPE.HEART,
        "cruz": MSO_SHAPE.CROSS,
        "seta_direita": MSO_SHAPE.RIGHT_ARROW,
        "seta_esquerda": MSO_SHAPE.LEFT_ARROW,
        "seta_cima": MSO_SHAPE.UP_ARROW,
        "seta_baixo": MSO_SHAPE.DOWN_ARROW,
        "chevron": MSO_SHAPE.CHEVRON,
        "documento": MSO_SHAPE.FLOWCHART_DOCUMENT,
        "processo": MSO_SHAPE.FLOWCHART_PROCESS,
        "decisao": MSO_SHAPE.FLOWCHART_DECISION,
        "dados": MSO_SHAPE.FLOWCHART_DATA,
        "cilindro": MSO_SHAPE.CAN,
        "moldura": MSO_SHAPE.FRAME,
        "donut": MSO_SHAPE.DONUT,
        "proibido": MSO_SHAPE.NO_SYMBOL,
        "arco": MSO_SHAPE.BLOCK_ARC,
        "retangulo": MSO_SHAPE.RECTANGLE,
        "retangulo_arredondado": MSO_SHAPE.ROUNDED_RECTANGLE,
    }

    def _adicionar_icone(self, slide, icone_cfg_raw, tema):
        """Adiciona um icone ao slide usando shapes nativos do PowerPoint.

        Aceita campos no formato SKILL.md (x, y, tamanho, etc.)
        ou no formato interno (left, top, size, etc.).
        """
        icone_cfg = _normalizar_cfg(icone_cfg_raw, _MAP_ICONE)

        left = Inches(icone_cfg.get("left", 0.5))
        top = Inches(icone_cfg.get("top", 0.5))
        size = Inches(icone_cfg.get("size", 0.5))
        width = Inches(icone_cfg.get("width", icone_cfg.get("size", 0.5)))
        height = Inches(icone_cfg.get("height", icone_cfg.get("size", 0.5)))

        # Resolver forma
        forma_nome = icone_cfg.get("forma", "circulo")
        shape_type = self.SHAPE_MAP.get(forma_nome, MSO_SHAPE.OVAL)

        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        # Fundo
        cor_fundo = icone_cfg.get("cor_fundo", "marinho")
        if cor_fundo and cor_fundo != "nenhum":
            self._set_shape_fill(shape, cor_fundo)
        else:
            shape.fill.background()

        # Borda
        borda = icone_cfg.get("borda", "nenhum")
        if borda and borda != "nenhum":
            self._set_shape_border(shape, borda, icone_cfg.get("borda_espessura", 1.5))
        else:
            shape.line.fill.background()

        # Rotacao
        if icone_cfg.get("rotacao"):
            shape.rotation = icone_cfg["rotacao"]

        # Texto/simbolo (opcional - para formas que suportam texto)
        simbolo = icone_cfg.get("simbolo", "")
        if simbolo:
            self._set_shape_text(
                shape, simbolo,
                font_name=icone_cfg.get("font_name", "Arial Nova Condensed"),
                font_size=icone_cfg.get("font_size", 14),
                color=icone_cfg.get("cor_texto", "branco"),
                bold=icone_cfg.get("bold", True),
                alignment="center",
                vertical="middle",
            )

    def _adicionar_callout(self, slide, callout_cfg_raw, tema):
        """Adiciona um callout de KPI (numero grande + label) ao slide.

        Aceita campos no formato SKILL.md (x, y, largura, valor, cor_valor, etc.)
        ou no formato interno (left, top, width, numero, cor_numero, etc.).
        """
        callout_cfg = _normalizar_cfg(callout_cfg_raw, _MAP_CALLOUT)

        # SKILL.md usa altura_valor + altura_label em vez de height unico
        if "height" not in callout_cfg:
            av = callout_cfg.get("altura_valor", 0.6)
            al = callout_cfg.get("altura_label", 0.4)
            af = callout_cfg.get("altura_fundo", None)
            callout_cfg["height"] = af if af else av + al

        left = Inches(callout_cfg.get("left", 0.5))
        top = Inches(callout_cfg.get("top", 2.0))
        width = Inches(callout_cfg.get("width", 2.0))
        height = Inches(callout_cfg.get("height", 1.2))

        # Fundo opcional (retangulo)
        fundo = callout_cfg.get("fundo", "nenhum")
        borda = callout_cfg.get("borda", "nenhum")
        if fundo != "nenhum" or borda != "nenhum":
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            if fundo != "nenhum":
                self._set_shape_fill(bg_shape, fundo)
            else:
                bg_shape.fill.background()
            if borda != "nenhum":
                self._set_shape_border(bg_shape, borda, callout_cfg.get("borda_espessura", 1.0))
            else:
                bg_shape.line.fill.background()

        # Numero grande
        numero = callout_cfg.get("numero", "")
        cor_num = self._resolver_cor(callout_cfg.get("cor_numero", "marinho"))
        num_size = callout_cfg.get("font_size_numero", 44)
        alignment = callout_cfg.get("alignment", "center")

        txBox_num = slide.shapes.add_textbox(left, top, width, Inches(height.inches * 0.6) if isinstance(height, Emu) else Inches(callout_cfg.get("height", 1.2) * 0.6))
        tf = txBox_num.text_frame
        tf.word_wrap = True
        tf.text = str(numero)
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        for p in tf.paragraphs:
            p.alignment = align_map.get(alignment, PP_ALIGN.CENTER)
            for run in p.runs:
                run.font.name = "Arial Nova Light"
                run.font.size = Pt(num_size)
                run.font.color.rgb = cor_num
                run.font.bold = True

        # Label
        label = callout_cfg.get("label", "")
        if label:
            cor_lbl = self._resolver_cor(callout_cfg.get("cor_label", "cinza_texto"))
            lbl_size = callout_cfg.get("font_size_label", 12)
            label_top_offset = callout_cfg.get("height", 1.2) * 0.55
            txBox_lbl = slide.shapes.add_textbox(left, Inches(callout_cfg.get("top", 2.0) + label_top_offset), width, Inches(callout_cfg.get("height", 1.2) * 0.45))
            tf2 = txBox_lbl.text_frame
            tf2.word_wrap = True
            tf2.text = str(label)
            for p in tf2.paragraphs:
                p.alignment = align_map.get(alignment, PP_ALIGN.CENTER)
                for run in p.runs:
                    run.font.name = "Arial Nova Condensed"
                    run.font.size = Pt(lbl_size)
                    run.font.color.rgb = cor_lbl

    def _adicionar_shape(self, slide, shape_cfg_raw, tema):
        """Adiciona um text box livre ao slide.

        Aceita campos no formato SKILL.md (x, y, largura, texto, etc.)
        ou no formato interno (left, top, width, text, etc.).
        """
        shape_cfg = _normalizar_cfg(shape_cfg_raw, _MAP_SHAPE)

        from pptx.util import Inches as In
        left = In(shape_cfg.get("left", 0.5))
        top = In(shape_cfg.get("top", 5.0))
        width = In(shape_cfg.get("width", 5.0))
        height = In(shape_cfg.get("height", 0.5))

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = shape_cfg.get("text", "")

        cor = shape_cfg.get("color", "cinza_texto")
        if isinstance(cor, str) and cor in CORES:
            cor = CORES[cor]
        else:
            cor = COR_TEXTO_POR_TEMA.get(tema, CORES["marinho"])

        for p in tf.paragraphs:
            for run in p.runs:
                run.font.name = shape_cfg.get("font_name", "Arial Nova Condensed")
                run.font.size = Pt(shape_cfg.get("font_size", 10))
                run.font.color.rgb = cor
                if shape_cfg.get("bold"):
                    run.font.bold = True


    # ====================================================================
    # COMPONENTES DD (extensao A&M): waterfall, quadrante 2x2, funil
    # ====================================================================
    def _adicionar_waterfall(self, slide, cfg_raw, tema):
        cfg = cfg_raw if isinstance(cfg_raw, dict) else {}
        left = float(cfg.get("x", cfg.get("left", 0.7)))
        top = float(cfg.get("y", cfg.get("top", 1.9)))
        width = float(cfg.get("largura", cfg.get("width", 11.9)))
        height = float(cfg.get("altura", cfg.get("height", 4.2)))
        barras = cfg.get("barras", [])
        if not barras:
            return
        cor_base = cfg.get("cor_base", "marinho")
        cor_up = cfg.get("cor_aumento", "azul_claro")
        cor_down = cfg.get("cor_reducao", "vermelho")
        fmt = cfg.get("formato", "{:.0f}")
        seq = []
        running = 0.0
        scale = [0.0]
        for b in barras:
            tipo_b = b.get("tipo", "aumento")
            v = float(b.get("valor", 0))
            if tipo_b in ("base", "total", "subtotal"):
                start, end = 0.0, v
                running = v
            elif tipo_b == "reducao":
                start, end = running, running - abs(v)
                running = end
            else:
                start, end = running, running + abs(v)
                running = end
            seq.append((tipo_b, v, start, end))
            scale.append(start); scale.append(end)
        vmax = max(scale); vmin = min(scale); rng = (vmax - vmin) or 1.0
        n = len(seq)
        gap_ratio = 0.4
        bar_w = width / (n * (1 + gap_ratio))
        plot_top = top + 0.35
        plot_h = height - 0.85
        def y_of(val):
            return plot_top + plot_h * (vmax - val) / rng
        for i, (tipo_b, v, start, end) in enumerate(seq):
            bx = left + i * bar_w * (1 + gap_ratio)
            y_hi = y_of(max(start, end)); y_lo = y_of(min(start, end))
            bh = max(y_lo - y_hi, 0.05)
            cor = cor_base if tipo_b in ("base", "total", "subtotal") else (cor_down if tipo_b == "reducao" else cor_up)
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(y_hi), Inches(bar_w), Inches(bh))
            self._set_shape_fill(shp, cor); shp.line.fill.background()
            lbl = slide.shapes.add_textbox(Inches(bx - 0.15), Inches(y_hi - 0.3), Inches(bar_w + 0.3), Inches(0.28))
            self._set_shape_text(lbl, fmt.format(v), font_size=9, color="marinho", bold=True, alignment="center", vertical="middle")
            cat = slide.shapes.add_textbox(Inches(bx - 0.2), Inches(plot_top + plot_h + 0.08), Inches(bar_w + 0.4), Inches(0.6))
            self._set_shape_text(cat, str(barras[i].get("label", "")), font_size=8, color="cinza_texto", bold=False, alignment="center", vertical="top")
            if i < n - 1 and tipo_b != "total":
                cy = y_of(end)
                conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx + bar_w), Inches(cy - 0.006), Inches(bar_w * gap_ratio), Inches(0.012))
                self._set_shape_fill(conn, "cinza_medio"); conn.line.fill.background()

    def _adicionar_quadrante(self, slide, cfg_raw, tema):
        cfg = cfg_raw if isinstance(cfg_raw, dict) else {}
        left = float(cfg.get("x", cfg.get("left", 1.4)))
        top = float(cfg.get("y", cfg.get("top", 1.9)))
        width = float(cfg.get("largura", cfg.get("width", 7.6)))
        height = float(cfg.get("altura", cfg.get("height", 4.2)))
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        frame.fill.background(); self._set_shape_border(frame, "cinza_medio", 1.0)
        midx = left + width / 2; midy = top + height / 2
        vbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(midx - 0.005), Inches(top), Inches(0.01), Inches(height))
        self._set_shape_fill(vbar, "cinza_medio"); vbar.line.fill.background()
        hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(midy - 0.005), Inches(width), Inches(0.01))
        self._set_shape_fill(hbar, "cinza_medio"); hbar.line.fill.background()
        quad = cfg.get("quadrantes", [])
        if len(quad) >= 4:
            pos = [(left + 0.12, top + 0.08), (midx + 0.12, top + 0.08), (left + 0.12, midy + 0.08), (midx + 0.12, midy + 0.08)]
            for (qx, qy), txt in zip(pos, quad):
                lb = slide.shapes.add_textbox(Inches(qx), Inches(qy), Inches(width / 2 - 0.24), Inches(0.3))
                self._set_shape_text(lb, str(txt), font_size=8, color="cinza_texto", bold=True, alignment="left", vertical="top")
        ex = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.06), Inches(width), Inches(0.3))
        self._set_shape_text(ex, cfg.get("eixo_x", ""), font_size=9, color="marinho", bold=True, alignment="center", vertical="middle")
        ey = slide.shapes.add_textbox(Inches(left - 0.1), Inches(top - 0.34), Inches(width), Inches(0.3))
        self._set_shape_text(ey, cfg.get("eixo_y", ""), font_size=9, color="marinho", bold=True, alignment="left", vertical="middle")
        for p in cfg.get("pontos", []):
            px = max(0.0, min(1.0, float(p.get("px", 0.5)))); py = max(0.0, min(1.0, float(p.get("py", 0.5))))
            d = float(p.get("tamanho", 0.45))
            cx = left + px * width; cy = top + (1 - py) * height
            ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
            self._set_shape_fill(ov, p.get("cor", "azul_claro")); ov.line.color.rgb = CORES["branco"]; ov.line.width = Pt(1.25)
            nm = slide.shapes.add_textbox(Inches(cx - 1.0), Inches(cy + d / 2 + 0.01), Inches(2.0), Inches(0.28))
            self._set_shape_text(nm, str(p.get("nome", "")), font_size=8, color="marinho", bold=False, alignment="center", vertical="top")

    def _adicionar_funil(self, slide, cfg_raw, tema):
        cfg = cfg_raw if isinstance(cfg_raw, dict) else {}
        left = float(cfg.get("x", cfg.get("left", 1.6)))
        top = float(cfg.get("y", cfg.get("top", 1.95)))
        width = float(cfg.get("largura", cfg.get("width", 6.2)))
        height = float(cfg.get("altura", cfg.get("height", 4.0)))
        niveis = cfg.get("niveis", [])
        if not niveis:
            return
        fmt = cfg.get("formato", "{}")
        n = len(niveis)
        gap = 0.14
        band_h = (height - (n - 1) * gap) / n
        grad = ["marinho", "cinza_inter", "azul_claro", "cinza_medio", "cinza_claro"]
        min_ratio = 0.42
        for i, nv in enumerate(niveis):
            w_i = width * (1 - (1 - min_ratio) * (i / max(n - 1, 1)))
            x_i = left + (width - w_i) / 2
            y_i = top + i * (band_h + gap)
            cor = nv.get("cor", grad[i % len(grad)])
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_i), Inches(y_i), Inches(w_i), Inches(band_h))
            self._set_shape_fill(band, cor); band.line.fill.background()
            cor_txt = "branco" if cor in ("marinho", "cinza_inter", "azul_escuro", "marinho_alt") else "marinho"
            self._set_shape_text(band, str(nv.get("label", "")), font_size=11, color=cor_txt, bold=True, alignment="center", vertical="middle")
            val = slide.shapes.add_textbox(Inches(left + width + 0.25), Inches(y_i + band_h / 2 - 0.2), Inches(3.4), Inches(0.5))
            vtxt = fmt.format(nv.get("valor", "")) if nv.get("valor", None) is not None else ""
            sub = nv.get("sublabel", "")
            full = vtxt + (("\n" + sub) if sub else "")
            self._set_shape_text(val, full, font_size=10, color="marinho", bold=True, alignment="left", vertical="middle")


# ============================================================================
# EXEMPLO DE USO
# ============================================================================

def _exemplo():
    """Retorna um JSON de exemplo para demonstracao."""
    return {
        "arquivo_saida": "Exemplo_AM.pptx",
        "slides": [
            # 1 - Capa
            {
                "tipo": "capa_grafismo_marinho",
                "titulo": "Projeto de\nTransformacao Digital",
                "subtitulo": "Nome do Cliente",
                "data": "Abril.2026"
            },
            # 2 - Agenda
            {
                "tipo": "agenda_azul",
                "titulo": "Projeto de\nTransformacao Digital",
                "itens": [
                    "Contexto e objetivos",
                    "Abordagem A&M",
                    "Metodologia",
                    "Time e investimento",
                    "Proximos passos"
                ]
            },
            # 3 - Divisoria
            {
                "tipo": "divisoria_marinho",
                "titulo": "Contexto e Objetivos"
            },
            # 4 - Conteudo padrao
            {
                "tipo": "padrao_branco",
                "tag": "Contexto",
                "titulo": "O cenario atual exige transformacao",
                "subtitulo": "Principais desafios identificados no diagnostico",
                "rodape": "Fonte: Analise A&M, 2026"
            },
            # 5 - Destaque
            {
                "tipo": "destaque_marinho",
                "texto": "A transformacao digital nao e sobre tecnologia.\nE sobre pessoas e processos.",
                "credito": "Visao A&M"
            },
            # 6 - Divisoria
            {
                "tipo": "divisoria_marinho",
                "titulo": "Abordagem A&M"
            },
            # 7 - 3 colunas
            {
                "tipo": "texto3_branco",
                "titulo": "Nossa abordagem em tres pilares",
                "tag": "Abordagem",
                "subtitulo": "Metodologia proprietaria A&M para transformacao",
                "colunas": [
                    {"titulo": "Diagnostico", "corpo": "Mapeamento completo de processos, sistemas e estrutura organizacional. Identificacao de gaps e oportunidades."},
                    {"titulo": "Planejamento", "corpo": "Definicao de roadmap com marcos claros, priorizacao de iniciativas e alocacao de recursos."},
                    {"titulo": "Execucao", "corpo": "Implementacao agil com acompanhamento semanal de KPIs e gestao de mudanca."}
                ],
                "rodape": "Nota: Abordagem adaptavel conforme necessidades do cliente"
            },
            # 8 - Conteudo com grafismo
            {
                "tipo": "grafismo_branco",
                "tag": "Metodologia",
                "titulo": "Framework de implementacao",
                "subtitulo": "Etapas do projeto com entregaveis definidos",
                "rodape": "Fonte: Metodologia A&M"
            },
            # 9 - 2 colunas
            {
                "tipo": "texto2_branco",
                "titulo": "Comparativo de cenarios",
                "tag": "Analise",
                "subtitulo": "Cenario atual versus cenario proposto",
                "colunas": [
                    {"titulo": "Cenario Atual", "corpo": "Processos manuais e fragmentados\nSilos de informacao entre areas\nAlto custo operacional\nBaixa visibilidade de dados"},
                    {"titulo": "Cenario Proposto", "corpo": "Processos automatizados e integrados\nPlataforma unica de dados\nReducao de 30% em custos\nDashboards em tempo real"}
                ]
            },
            # 10 - Topicos
            {
                "tipo": "topicos_claro",
                "titulo": "Principais entregas do projeto",
                "tag": "Entregas",
                "subtitulo": "Resultados esperados por fase",
                "destaque": "O projeto sera conduzido em 3 fases com duracaototal estimada de 16 semanas, entregando valor incremental desde a primeira semana.",
                "topicos": [
                    {"titulo": "Fase 1", "corpo": "Diagnostico e quick wins (4 semanas)"},
                    {"titulo": "Fase 2", "corpo": "Implementacao dos pilares prioritarios (8 semanas)"},
                    {"titulo": "Fase 3", "corpo": "Estabilizacao e handover (4 semanas)"},
                    {"titulo": "Governanca", "corpo": "Comite semanal de acompanhamento"},
                    {"titulo": "KPIs", "corpo": "Metricas de sucesso definidas por fase"}
                ]
            },
            # 11 - Divisoria
            {
                "tipo": "divisoria_marinho",
                "titulo": "Time e Investimento"
            },
            # 12 - Texto simples
            {
                "tipo": "texto1_branco",
                "titulo": "Proposta de investimento",
                "tag": "Comercial",
                "subtitulo": "Valores e condicoes",
                "col1_titulo": "Estrutura de fees",
                "col1_corpo": "O investimento sera estruturado em fees mensais fixos durante o periodo do projeto, com possibilidade de success fee atrelado aos resultados alcancados.\n\nOs valores contemplam dedicacao do time senior, ferramentas proprietarias e metodologia A&M.",
                "rodape": "Nota: Valores sujeitos a aprovacao comercial"
            },
            # 13 - Despedida
            {
                "tipo": "despedida"
            },
        ]
    }


# ============================================================================
# MAIN
# ============================================================================

def main():
    # Determinar diretorio do script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "Template_Base.pptx")

    if not os.path.exists(template_path):
        print(f"ERRO: Template_Base.pptx nao encontrado em {script_dir}")
        print("Execute primeiro o script de criacao do template base.")
        sys.exit(1)

    # Carregar configuracao
    if len(sys.argv) > 1:
        json_path = sys.argv[1]
        if not os.path.exists(json_path):
            print(f"ERRO: Arquivo {json_path} nao encontrado.")
            sys.exit(1)
        with open(json_path, "r", encoding="utf-8") as f:
            config = json.load(f)
        print(f"Carregado: {json_path}")
    else:
        print("Nenhum JSON informado. Usando exemplo de demonstracao.")
        config = _exemplo()

    # Gerar
    output = config.get("arquivo_saida", "Apresentacao.pptx")
    output_path = os.path.join(script_dir, output)
    config["arquivo_saida"] = output_path

    gerador = GeradorAM(template_path)
    gerador.gerar(config)


if __name__ == "__main__":
    main()
